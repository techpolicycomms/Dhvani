#!/usr/bin/env python3
"""
Build the Dhvani product-development + roadmap pitch deck as a .pptx.

Single-file generator. Zero external network calls. Output goes to
docs/presentation/Dhvani_Roadmap.pptx. Re-run any time the plan changes.

Usage:
    /tmp/dhvani-pptx-venv/bin/python scripts/build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ---------------------------------------------------------------------
# Brand tokens (match lib/theme.ts + app/globals.css).
# ---------------------------------------------------------------------
ITU_BLUE = RGBColor(0x00, 0x9C, 0xD6)
ITU_BLUE_DARK = RGBColor(0x00, 0x7A, 0xAD)
ITU_BLUE_LIGHT = RGBColor(0x4F, 0xB7, 0xE4)
ITU_BLUE_PALE = RGBColor(0xE5, 0xF4, 0xFB)
DARK_NAVY = RGBColor(0x00, 0x33, 0x66)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)
BORDER_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x05, 0x96, 0x69)
WARNING = RGBColor(0xD9, 0x77, 0x06)

# Shared layout constants (16:9 at 13.333" x 7.5").
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Side rails.
CONTENT_LEFT = Inches(0.6)
CONTENT_TOP = Inches(1.4)
CONTENT_W = Inches(12.13)
CONTENT_H = Inches(5.4)


# ---------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    """Filled rectangle. Optional visible outline."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = fill_color
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    font_size=14,
    bold=False,
    color=DARK_GRAY,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    anchor=MSO_ANCHOR.TOP,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_brand_rail(slide):
    """Thin ITU-blue rail across the top of the slide."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ITU_BLUE)


def add_slide_header(slide, eyebrow, title):
    """Slide title block: small eyebrow + big title."""
    add_brand_rail(slide)
    add_text(
        slide,
        eyebrow,
        CONTENT_LEFT,
        Inches(0.35),
        CONTENT_W,
        Inches(0.3),
        font_size=11,
        bold=True,
        color=ITU_BLUE_DARK,
    )
    add_text(
        slide,
        title,
        CONTENT_LEFT,
        Inches(0.62),
        CONTENT_W,
        Inches(0.7),
        font_size=28,
        bold=True,
        color=DARK_NAVY,
    )


def add_footer(slide, slide_no, total):
    add_text(
        slide,
        "Dhvani  ·  Internal  ·  2026",
        CONTENT_LEFT,
        Inches(7.1),
        Inches(6),
        Inches(0.3),
        font_size=9,
        color=MID_GRAY,
    )
    add_text(
        slide,
        f"{slide_no} / {total}",
        Inches(12.2),
        Inches(7.1),
        Inches(1),
        Inches(0.3),
        font_size=9,
        color=MID_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, bullets, left, top, width, height, font_size=16, gap=0.08):
    """
    bullets: list[str] or list[tuple(level, str)]. Level 0 is primary, 1 is
    a sub-bullet with smaller text.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6 if level == 0 else 3)
        marker = "•" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{marker}  {text}"
        run.font.size = Pt(font_size if level == 0 else font_size - 2)
        run.font.color.rgb = DARK_GRAY if level == 0 else MID_GRAY
        run.font.name = "Calibri"
        if level == 0:
            run.font.bold = False


def add_pill(slide, text, left, top, width, height, fill, text_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"


# ---------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------
def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Left-edge ITU-blue panel.
    add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, ITU_BLUE)
    # Thin rule at top.
    add_rect(slide, Inches(0.3), Inches(0.7), Inches(1.2), Inches(0.06), ITU_BLUE)
    add_text(
        slide,
        "DHVANI",
        Inches(0.9),
        Inches(1.2),
        Inches(12),
        Inches(1),
        font_size=56,
        bold=True,
        color=DARK_NAVY,
    )
    add_text(
        slide,
        "Product Development & 12-Month Roadmap",
        Inches(0.9),
        Inches(2.3),
        Inches(12),
        Inches(0.8),
        font_size=28,
        color=ITU_BLUE_DARK,
    )
    add_text(
        slide,
        "From Innovation Pilot to ITU's Default Transcription Tool",
        Inches(0.9),
        Inches(3.1),
        Inches(12),
        Inches(0.6),
        font_size=18,
        color=DARK_GRAY,
    )
    # Metadata line.
    add_text(
        slide,
        "Presented to   Director  ·  CIO  ·  ISD Leadership",
        Inches(0.9),
        Inches(5.5),
        Inches(12),
        Inches(0.4),
        font_size=14,
        bold=True,
        color=DARK_NAVY,
    )
    add_text(
        slide,
        "ITU Innovation Hub   ·   2026",
        Inches(0.9),
        Inches(5.9),
        Inches(12),
        Inches(0.4),
        font_size=12,
        color=MID_GRAY,
    )
    # Pill: "Internal — confidential"
    add_pill(
        slide,
        "INTERNAL  ·  CONFIDENTIAL",
        Inches(0.9),
        Inches(6.6),
        Inches(2.3),
        Inches(0.35),
        ITU_BLUE_PALE,
        ITU_BLUE_DARK,
    )


def build_claim_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "THE CLAIM", "ITU does not need to buy a transcription tool.")
    # Main body paragraph.
    add_text(
        slide,
        "ITU needs a tool native to how ITU actually works.",
        CONTENT_LEFT,
        CONTENT_TOP,
        CONTENT_W,
        Inches(0.6),
        font_size=22,
        bold=True,
        color=DARK_NAVY,
    )
    bullets = [
        "Six UN languages, with mid-sentence code-switch",
        "Study Group and Bureau acronyms built in by default",
        "Delegates attributed by country and role",
        "Formal ITU-style record outputs — not plain transcripts",
        "Data inside the ITU Azure tenant, never outside",
        "Under $2 per user per month at scale",
    ]
    add_bullets(slide, bullets, CONTENT_LEFT, Inches(2.3), CONTENT_W, Inches(3), font_size=18)
    # Closing line.
    add_rect(slide, CONTENT_LEFT, Inches(5.7), CONTENT_W, Inches(0.04), ITU_BLUE_LIGHT)
    add_text(
        slide,
        "Dhvani already has the engine. In 12 months, Dhvani becomes the default tool "
        "every ITU staff member, delegate, and observer reaches for.",
        CONTENT_LEFT,
        Inches(5.9),
        CONTENT_W,
        Inches(1),
        font_size=14,
        color=DARK_GRAY,
    )
    add_footer(slide, n, total)


def build_problem_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "PROBLEM TODAY", "ITU staff lose hours each week.")
    bullets = [
        "Re-listening to meetings to extract what was actually decided",
        "Translating quotes across six official UN languages",
        "Hand-typing delegate attributions and country names into notes",
        "Rewriting transcripts into ITU's formal record formats",
        "Searching across Outlook, OneDrive, and notebooks for \"what we decided last time\"",
    ]
    add_bullets(slide, bullets, CONTENT_LEFT, CONTENT_TOP, CONTENT_W, Inches(4), font_size=20)

    # Emphasis callout.
    callout_top = Inches(6.0)
    add_rect(slide, CONTENT_LEFT, callout_top, CONTENT_W, Inches(0.7), ITU_BLUE_PALE)
    add_text(
        slide,
        "Generic transcription tools don't solve any of this. They weren't designed for the UN.",
        CONTENT_LEFT + Inches(0.3),
        callout_top + Inches(0.15),
        CONTENT_W - Inches(0.6),
        Inches(0.5),
        font_size=15,
        bold=True,
        color=ITU_BLUE_DARK,
    )
    add_footer(slide, n, total)


def build_commercial_table_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "WHY COMMERCIAL TOOLS FAIL ITU",
        "One table — for any Bureau director who asks.",
    )
    rows = [
        (
            "Six UN languages + mid-sentence code-switch",
            "Weak or English-only",
            "First-class",
        ),
        (
            "ITU vocabulary (SG-17, WTSA, CPM, RRB, PP, TD)",
            "Treated as noise",
            "Pre-loaded default",
        ),
        (
            "Delegate recognition by country and role",
            "Generic \"Speaker 1\"",
            "\"Canada (Chair)\", \"Secretariat\"",
        ),
        (
            "Formal output templates (ITU-T, ITU-R, verbatim)",
            "Plain transcript only",
            "Native ITU templates",
        ),
        (
            "Data residency in ITU Azure tenant",
            "No",
            "Yes",
        ),
        (
            "Per-user cost at 10,000 staff",
            "Prohibitive per-seat",
            "Under $2/month",
        ),
    ]
    headers = ["Capability", "Commercial SaaS", "Dhvani (planned)"]
    cols = [Inches(5.5), Inches(3.4), Inches(3.3)]
    start_left = CONTENT_LEFT
    top = Inches(1.55)
    row_h = Inches(0.55)

    # Header row.
    left = start_left
    for i, h in enumerate(headers):
        add_rect(slide, left, top, cols[i], row_h, DARK_NAVY)
        add_text(
            slide,
            h,
            left + Inches(0.15),
            top + Inches(0.1),
            cols[i] - Inches(0.3),
            row_h,
            font_size=12,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        left += cols[i]

    # Body rows.
    top += row_h
    for idx, r in enumerate(rows):
        left = start_left
        bg = WHITE if idx % 2 == 0 else OFF_WHITE
        for i, cell in enumerate(r):
            add_rect(slide, left, top, cols[i], row_h, bg)
            add_text(
                slide,
                cell,
                left + Inches(0.15),
                top + Inches(0.1),
                cols[i] - Inches(0.3),
                row_h,
                font_size=11,
                color=DARK_GRAY if i != 2 else ITU_BLUE_DARK,
                bold=(i == 2),
                anchor=MSO_ANCHOR.MIDDLE,
            )
            left += cols[i]
        top += row_h
    add_footer(slide, n, total)


def build_status_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "WHERE DHVANI IS TODAY", "Ninety percent of the foundation is done.")
    bullets = [
        "Next.js 14 PWA + Electron desktop wrapper",
        "Microsoft Entra SSO inside the ITU tenant",
        "Azure OpenAI transcription (gpt-4o-transcribe-diarize) + summary (gpt-4.1-mini), server-side only",
        "Personal and Power modes for different user workflows",
        "Dark mode, keyboard shortcuts, WCAG-AA palette",
        "Crash recovery with OPFS chunk persistence + silent auto-retry",
        "Azure Blob storage backend, durable across redeploys",
        "Official ITU brand identity — #009CD6 blue, Noto Sans, lucide iconography",
    ]
    add_bullets(slide, bullets, CONTENT_LEFT, CONTENT_TOP, CONTENT_W, Inches(5), font_size=17)
    add_rect(slide, CONTENT_LEFT, Inches(6.4), CONTENT_W, Inches(0.04), ITU_BLUE_LIGHT)
    add_text(
        slide,
        "Ready for internal beta.",
        CONTENT_LEFT,
        Inches(6.5),
        CONTENT_W,
        Inches(0.5),
        font_size=16,
        bold=True,
        color=ITU_BLUE_DARK,
    )
    add_footer(slide, n, total)


def build_roadmap_overview_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "THE 12-MONTH ROADMAP AT A GLANCE",
        "Four phases. One goal each. Explicit non-goals.",
    )
    phases = [
        ("Phase 1", "Weeks 0–4", "Polish & Trust", "Ship to 50 HQ beta users", ITU_BLUE_PALE, ITU_BLUE_DARK),
        ("Phase 2", "Weeks 5–12", "ITU-Native", "\"Oh, this is made for us\"", ITU_BLUE_LIGHT, DARK_NAVY),
        ("Phase 3", "Weeks 13–26", "Workflow Integration", "Feels like part of ITU", ITU_BLUE, WHITE),
        ("Phase 4", "Weeks 27–52", "Intelligence", "Post-meeting value > meeting value", ITU_BLUE_DARK, WHITE),
    ]
    card_w = Inches(2.9)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start_left = (SLIDE_W - total_w) / 2
    top = Inches(2.0)
    for i, (tag, weeks, title, goal, fill, txt) in enumerate(phases):
        left = start_left + (card_w + gap) * i
        add_rect(slide, left, top, card_w, Inches(3.4), fill)
        add_text(
            slide,
            tag,
            left + Inches(0.25),
            top + Inches(0.25),
            card_w - Inches(0.5),
            Inches(0.35),
            font_size=11,
            bold=True,
            color=txt,
        )
        add_text(
            slide,
            weeks,
            left + Inches(0.25),
            top + Inches(0.55),
            card_w - Inches(0.5),
            Inches(0.3),
            font_size=10,
            color=txt,
        )
        add_text(
            slide,
            title,
            left + Inches(0.25),
            top + Inches(1.1),
            card_w - Inches(0.5),
            Inches(1),
            font_size=22,
            bold=True,
            color=txt,
        )
        add_text(
            slide,
            goal,
            left + Inches(0.25),
            top + Inches(2.2),
            card_w - Inches(0.5),
            Inches(1.1),
            font_size=13,
            color=txt,
        )
    add_text(
        slide,
        "Non-goals are explicit at every phase. We do not ship features not on the list.",
        CONTENT_LEFT,
        Inches(6.0),
        CONTENT_W,
        Inches(0.4),
        font_size=12,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, n, total)


def build_phase_slide(prs, n, total, phase_no, weeks, title, goal, ships, metric):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, f"PHASE {phase_no}  ·  {weeks}", title)
    # Goal box.
    add_rect(slide, CONTENT_LEFT, Inches(1.55), CONTENT_W, Inches(0.6), ITU_BLUE_PALE)
    add_text(
        slide,
        f"Goal:  {goal}",
        CONTENT_LEFT + Inches(0.25),
        Inches(1.65),
        CONTENT_W - Inches(0.5),
        Inches(0.4),
        font_size=14,
        bold=True,
        color=ITU_BLUE_DARK,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Must-ship bullets.
    add_text(
        slide,
        "Must ship",
        CONTENT_LEFT,
        Inches(2.4),
        CONTENT_W,
        Inches(0.4),
        font_size=13,
        bold=True,
        color=DARK_NAVY,
    )
    add_bullets(slide, ships, CONTENT_LEFT, Inches(2.85), CONTENT_W, Inches(3.5), font_size=14)
    # Metric gate.
    add_rect(slide, CONTENT_LEFT, Inches(6.3), CONTENT_W, Inches(0.04), ITU_BLUE_LIGHT)
    add_text(
        slide,
        f"Metric gate:  {metric}",
        CONTENT_LEFT,
        Inches(6.4),
        CONTENT_W,
        Inches(0.5),
        font_size=13,
        bold=True,
        color=ITU_BLUE_DARK,
    )
    add_footer(slide, n, total)


def build_demo_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "THE \"WOW\" DEMO",
        "Five minutes. Memorize it. No commercial tool does this today.",
    )
    steps = [
        ("09:58", "Open Outlook. Click \"Transcribe\" on a Teams invite titled SG-17 Rapporteur · Security."),
        ("09:58", "Dhvani opens pre-configured. Banner: \"Ready for SG-17 · starts in 2 min.\""),
        ("10:00", "Meeting begins. Transcript fills live. Speakers labeled Kenya, France (Chair), Secretariat."),
        ("10:04", "Delegate switches to French mid-sentence. Transcript captures French faithfully; English gloss toggle."),
        ("10:47", "Meeting ends. Click \"Generate ITU-T summary record\". Formal .docx with ITU header + TOC."),
        ("17:00", "Email digest: \"From today's SG-17 meeting: 3 decisions, 4 action items, 2 are yours.\""),
        ("Later", "Cmd+K: \"what did we decide about spectrum for IMT-2030?\" — cited answer across 3 meetings."),
    ]
    top = Inches(1.55)
    row_h = Inches(0.72)
    time_w = Inches(1.3)
    text_w = CONTENT_W - time_w - Inches(0.3)
    for i, (time, text) in enumerate(steps):
        # Time pill.
        add_pill(
            slide,
            time,
            CONTENT_LEFT,
            top + Inches(0.12),
            time_w,
            Inches(0.45),
            ITU_BLUE,
            WHITE,
        )
        add_text(
            slide,
            text,
            CONTENT_LEFT + time_w + Inches(0.3),
            top + Inches(0.15),
            text_w,
            row_h,
            font_size=14,
            color=DARK_GRAY,
            anchor=MSO_ANCHOR.TOP,
        )
        # Connector line except on last row.
        if i < len(steps) - 1:
            add_rect(
                slide,
                CONTENT_LEFT + time_w / 2 - Emu(4500),
                top + Inches(0.58),
                Emu(9000),
                row_h - Inches(0.12),
                ITU_BLUE_LIGHT,
            )
        top += row_h
    add_footer(slide, n, total)


def build_ui_polish_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "UI/UX POLISH",
        "Five visible shifts that change \"pilot\" to \"product\".",
    )
    items = [
        ("Onboarding", "60-second welcome flow that ends with confetti and a transcript — not a docs link."),
        ("Record button", "Big circular breathing button with live waveform + speaker-color highlighting."),
        ("Speaker pills", "Country flag + role: 🇰🇪 Kenya, 🇫🇷 France · Chair, Secretariat."),
        ("Recap preview", "Formal-document view that looks like an ITU TD — header, TOC, page numbers."),
        ("Command palette", "Cmd+K — the power-user signal that makes engineers take you seriously."),
    ]
    left = CONTENT_LEFT
    top = Inches(1.6)
    card_w = Inches(3.9)
    card_h = Inches(2.2)
    gap = Inches(0.2)
    for i, (label, text) in enumerate(items):
        col = i % 3
        row = i // 3
        x = left + (card_w + gap) * col
        y = top + (card_h + gap) * row
        add_rect(slide, x, y, card_w, card_h, OFF_WHITE)
        add_rect(slide, x, y, Inches(0.08), card_h, ITU_BLUE)
        add_text(
            slide,
            label,
            x + Inches(0.3),
            y + Inches(0.2),
            card_w - Inches(0.4),
            Inches(0.4),
            font_size=13,
            bold=True,
            color=ITU_BLUE_DARK,
        )
        add_text(
            slide,
            text,
            x + Inches(0.3),
            y + Inches(0.7),
            card_w - Inches(0.4),
            card_h - Inches(0.9),
            font_size=13,
            color=DARK_GRAY,
        )
    add_text(
        slide,
        "Plus: Noto Sans with tabular numerals  ·  ITU blue used sparingly as accent  ·  dark mode as designed, not inverted.",
        CONTENT_LEFT,
        Inches(6.4),
        CONTENT_W,
        Inches(0.4),
        font_size=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, n, total)


def build_priority_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "FEATURE PRIORITIZATION",
        "Ruthless scope control. Nothing off-list ships.",
    )
    cols = [
        ("★★★★★  Five-star — must do", ITU_BLUE, [
            "Multilingual code-switch (P2)",
            "ITU vocabulary pack (P2, small)",
            "Delegate recognition (P2)",
        ]),
        ("★★★★  Four-star", ITU_BLUE_DARK, [
            "Formal export templates (P2)",
            "Outlook add-in (P3)",
            "Teams bot (P3, needs policy)",
            "Ask Dhvani personal RAG (P4)",
            "Live translation view (P4)",
        ]),
        ("★★★  Three-star", MID_GRAY, [
            "Library search + fuzzy filter (P2)",
            "Command palette / Cmd+K (P2)",
            "Mobile-native polish (P3)",
            "SharePoint sync (P3)",
            "Email digest (P3)",
            "Bureau analytics (P4)",
        ]),
    ]
    left = CONTENT_LEFT
    top = Inches(1.55)
    card_w = Inches(3.9)
    card_h = Inches(5.0)
    gap = Inches(0.2)
    for i, (title, color, items) in enumerate(cols):
        x = left + (card_w + gap) * i
        add_rect(slide, x, top, card_w, Inches(0.6), color)
        add_text(
            slide,
            title,
            x + Inches(0.2),
            top + Inches(0.1),
            card_w - Inches(0.4),
            Inches(0.4),
            font_size=12,
            bold=True,
            color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_rect(slide, x, top + Inches(0.6), card_w, card_h - Inches(0.6), OFF_WHITE)
        add_bullets(
            slide,
            items,
            x + Inches(0.2),
            top + Inches(0.8),
            card_w - Inches(0.4),
            card_h - Inches(0.9),
            font_size=12,
        )
    add_footer(slide, n, total)


def build_metrics_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "SUCCESS METRICS AT 12 MONTHS",
        "Shared monthly with Bureau directors.",
    )
    metrics = [
        ("3,000", "Monthly active users across HQ + regional offices"),
        (">50%", "of internal Teams meetings transcribed via Dhvani"),
        ("<5%", "mis-labeled-speaker rate on the ITU eval set"),
        ("<$2", "per user per month all-in"),
        ("99.5%", "availability; p95 chunk latency < 6s"),
        (">90%", "code-switch segments transcribed correctly"),
        ("1+", "external UN agency formally requesting access"),
        ("3+", "unsolicited mentions in a Bureau director meeting"),
    ]
    left = CONTENT_LEFT
    top = Inches(1.6)
    card_w = Inches(2.9)
    card_h = Inches(2.3)
    gap_x = Inches(0.18)
    gap_y = Inches(0.18)
    for i, (big, label) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = left + (card_w + gap_x) * col
        y = top + (card_h + gap_y) * row
        add_rect(slide, x, y, card_w, card_h, OFF_WHITE)
        add_rect(slide, x, y, card_w, Inches(0.08), ITU_BLUE)
        add_text(
            slide,
            big,
            x + Inches(0.2),
            y + Inches(0.4),
            card_w - Inches(0.4),
            Inches(1),
            font_size=40,
            bold=True,
            color=ITU_BLUE_DARK,
        )
        add_text(
            slide,
            label,
            x + Inches(0.2),
            y + Inches(1.45),
            card_w - Inches(0.4),
            card_h - Inches(1.5),
            font_size=11,
            color=DARK_GRAY,
        )
    add_footer(slide, n, total)


def build_risks_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "RISKS  &  HONEST MITIGATIONS",
        "We name the risks. We name the fixes. No hand-waving.",
    )
    rows = [
        ("Multilingual quality plateau", "Dedicated eng-month in P2 on per-language prompt tuning + eval harness."),
        ("Teams bot blocked by M365 policy", "Raise to CIO at month 3. Product works without it."),
        ("LLM cost drift at scale", "Per-user + monthly org caps already live. PTU at 10k users."),
        ("Privacy pushback on voice archival", "Archival is OFF by default. Turn on only with legal sign-off."),
        ("Adoption plateau", "Champions program, leadership seeding, office hours, peer-pressure dashboard."),
        ("Bus factor", "Two engineers — not one. HANDOFF.md + runbooks in place."),
    ]
    top = Inches(1.6)
    row_h = Inches(0.75)
    left_w = Inches(5.0)
    right_w = CONTENT_W - left_w - Inches(0.2)
    for i, (risk, mit) in enumerate(rows):
        y = top + row_h * i
        bg = WHITE if i % 2 == 0 else OFF_WHITE
        add_rect(slide, CONTENT_LEFT, y, CONTENT_W, row_h, bg)
        add_rect(slide, CONTENT_LEFT, y, Inches(0.08), row_h, WARNING)
        add_text(
            slide,
            risk,
            CONTENT_LEFT + Inches(0.25),
            y + Inches(0.12),
            left_w - Inches(0.3),
            row_h,
            font_size=13,
            bold=True,
            color=DARK_NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            mit,
            CONTENT_LEFT + left_w,
            y + Inches(0.12),
            right_w,
            row_h,
            font_size=12,
            color=DARK_GRAY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    add_footer(slide, n, total)


def build_ask_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "THE ASK  —  FIVE UNBLOCKS",
        "Five decisions from leadership make this plan real.",
    )
    asks = [
        ("1", "Budget", "2 engineers + 0.5 designer × 9 months  ·  ≈ $400k all-in  ·  cheaper than buying commercial at ITU scale."),
        ("2", "Azure OpenAI TPM", "Raise tokens-per-minute to 500K per region."),
        ("3", "Entra registration", "Migrate to an ISD-owned app registry so we can issue production-grade tokens."),
        ("4", "Privacy pathway", "Open a review track for voice-audio archival; it is scaffolded but not wired."),
        ("5", "Executive sponsor", "Name one Bureau director as champion and first live-adopter."),
    ]
    top = Inches(1.6)
    row_h = Inches(0.95)
    num_w = Inches(0.7)
    label_w = Inches(2.5)
    for i, (num, label, text) in enumerate(asks):
        y = top + row_h * i
        add_rect(slide, CONTENT_LEFT, y, num_w, row_h - Inches(0.1), ITU_BLUE)
        add_text(
            slide,
            num,
            CONTENT_LEFT,
            y + Inches(0.1),
            num_w,
            row_h - Inches(0.3),
            font_size=28,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            label,
            CONTENT_LEFT + num_w + Inches(0.3),
            y + Inches(0.15),
            label_w,
            row_h - Inches(0.3),
            font_size=16,
            bold=True,
            color=DARK_NAVY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            text,
            CONTENT_LEFT + num_w + label_w + Inches(0.5),
            y + Inches(0.15),
            CONTENT_W - num_w - label_w - Inches(0.5),
            row_h - Inches(0.3),
            font_size=13,
            color=DARK_GRAY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    add_text(
        slide,
        "Say yes to these five. We commit to the plan. Everything else we'll figure out.",
        CONTENT_LEFT,
        Inches(6.6),
        CONTENT_W,
        Inches(0.4),
        font_size=13,
        bold=True,
        color=ITU_BLUE_DARK,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, n, total)


def build_next14_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "NEXT 14 DAYS  —  NO APPROVALS NEEDED",
        "Work that starts Monday. No blockers.",
    )
    items = [
        "Land current branch to main after PR review",
        "Stand up App Insights; wire five north-star metrics",
        "Ship ITU vocabulary pack v1 (150 hand-curated acronyms)",
        "Add Cmd+K command palette (hooks into existing shortcuts)",
        "Run first E2E test suite against the build; publish the report",
        "Produce the 5-slide demo deck of the \"wow\" moment",
        "Identify five executive beta users; send invite",
        "Draft the Phase 2 engineering plan; present at next staff meeting",
    ]
    # 2-column checklist layout.
    left_col = items[: len(items) // 2 + len(items) % 2]
    right_col = items[len(items) // 2 + len(items) % 2 :]
    for col_idx, col in enumerate([left_col, right_col]):
        x = CONTENT_LEFT + Inches(6.1 * col_idx)
        top = Inches(1.7)
        for i, text in enumerate(col):
            y = top + Inches(0.6 * i)
            # Checkbox square.
            add_rect(slide, x, y + Inches(0.08), Inches(0.22), Inches(0.22), WHITE, line=True)
            checkbox = slide.shapes[-1]
            checkbox.line.color.rgb = ITU_BLUE
            add_text(
                slide,
                text,
                x + Inches(0.4),
                y,
                Inches(5.5),
                Inches(0.5),
                font_size=13,
                color=DARK_GRAY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
    add_footer(slide, n, total)


def build_closing_slide(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ITU_BLUE)
    add_text(
        slide,
        "WHY NOW",
        CONTENT_LEFT,
        Inches(0.5),
        CONTENT_W,
        Inches(0.5),
        font_size=14,
        bold=True,
        color=ITU_BLUE_LIGHT,
    )
    add_text(
        slide,
        "Three windows are open simultaneously.",
        CONTENT_LEFT,
        Inches(1.1),
        CONTENT_W,
        Inches(1),
        font_size=36,
        bold=True,
        color=WHITE,
    )
    points = [
        "Azure OpenAI capability has caught up to ITU's multilingual needs (it wasn't true a year ago).",
        "Commercial transcription is consolidating on English-first, general-purpose — the UN niche is structurally unserved.",
        "ITU's Innovation Hub has built and proven the engine. The graduation window to ISD is open now.",
    ]
    top = Inches(2.6)
    for p in points:
        add_rect(slide, CONTENT_LEFT, top + Inches(0.22), Inches(0.2), Inches(0.2), ITU_BLUE)
        add_text(
            slide,
            p,
            CONTENT_LEFT + Inches(0.45),
            top,
            CONTENT_W - Inches(0.45),
            Inches(0.8),
            font_size=16,
            color=WHITE,
            anchor=MSO_ANCHOR.TOP,
        )
        top += Inches(0.8)

    add_rect(slide, CONTENT_LEFT, Inches(5.6), CONTENT_W, Inches(0.04), ITU_BLUE_LIGHT)
    add_text(
        slide,
        "Move in the next twelve months.  Dhvani becomes ITU's tool.",
        CONTENT_LEFT,
        Inches(5.85),
        CONTENT_W,
        Inches(0.6),
        font_size=22,
        bold=True,
        color=WHITE,
    )
    add_text(
        slide,
        "Wait, and someone procures a commercial alternative — and we spend three years working around its limitations.",
        CONTENT_LEFT,
        Inches(6.5),
        CONTENT_W,
        Inches(0.6),
        font_size=14,
        color=ITU_BLUE_LIGHT,
    )
    add_text(
        slide,
        "Thank you.  Questions welcome.",
        CONTENT_LEFT,
        Inches(7.05),
        CONTENT_W,
        Inches(0.3),
        font_size=11,
        color=MID_GRAY,
    )


# ---------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Total slide count for footer numbering (18 slides).
    total = 18

    build_title_slide(prs)  # 1
    build_claim_slide(prs, 2, total)  # 2
    build_problem_slide(prs, 3, total)  # 3
    build_commercial_table_slide(prs, 4, total)  # 4
    build_status_slide(prs, 5, total)  # 5
    build_roadmap_overview_slide(prs, 6, total)  # 6

    # Phases 1-4
    build_phase_slide(
        prs,
        7,
        total,
        phase_no=1,
        weeks="Weeks 0–4",
        title="Polish & Trust",
        goal="50 HQ beta users, no embarrassing moments, ready to invite leadership.",
        ships=[
            "Merge current branch to main; tag v1.0.0-rc1",
            "Azure Blob transcript storage ON in production",
            "Entra + secrets migrated to Azure Key Vault",
            "App Insights instrumentation with 5 north-star metrics",
            "GitHub Actions release pipeline — signed .dmg and .exe per tag",
            "WCAG 2.2 AA audit + one screen-reader walkthrough",
            "60-second onboarding flow that lands the first transcript",
        ],
        metric="50 beta users, no P1 incident, NPS > 30, ISD sign-off on v1.",
    )
    build_phase_slide(
        prs,
        8,
        total,
        phase_no=2,
        weeks="Weeks 5–12",
        title="ITU-Native",
        goal="The first demo where an ITU person says \"this is actually made for us\".",
        ships=[
            "Multilingual robustness — per-chunk language hint, mid-sentence code-switch, RTL Arabic tuned",
            "ITU vocabulary pack — 150+ acronyms pre-loaded by default (SG, WG, CPM, PP, RRB, TD)",
            "Delegate recognition — auto-pull attendee list from Outlook invite, map to Country + Role",
            "Formal export templates — ITU-T summary, ITU-R CPM, plenary verbatim, Bureau memo",
            "Meeting directory — saved transcripts auto-tagged with Study Group, Question, Bureau",
        ],
        metric="500 active users. Two Study Groups as primary adopters. One director quote.",
    )
    build_phase_slide(
        prs,
        9,
        total,
        phase_no=3,
        weeks="Weeks 13–26",
        title="Workflow Integration",
        goal="Dhvani feels like part of ITU, not a separate tool bolted on.",
        ships=[
            "Outlook add-in — \"Transcribe this meeting\" button on calendar invites",
            "Teams bot — joins as passive participant (with CIO policy approval)",
            "SharePoint sync — transcripts + action items land in Bureau libraries",
            "Planner / Tasks — extracted action items become assigned tasks",
            "Email digest — yesterday's meetings, action items, decisions owed",
            "Document pipeline — export into ITU-T/R/D document systems as draft TDs",
            "Mobile-native polish — first-class phone experience",
        ],
        metric="2,000 active users. 30% of internal Teams meetings transcribed. First external UN agency requests access.",
    )
    build_phase_slide(
        prs,
        10,
        total,
        phase_no=4,
        weeks="Weeks 27–52",
        title="Intelligence",
        goal="Post-meeting value exceeds the meeting value.",
        ships=[
            "Library search — fuzzy search + filter chips across everything you've recorded",
            "Ask Dhvani — personal RAG over your meeting history, with citations",
            "Meeting analytics per Bureau — talk-time equity, decision velocity, topic trends",
            "Live translation view — French speaker → English subtitles during live recording",
            "Automatic decision log per Study Group, curated from all its transcripts",
            "Live captions in the recording view; sign-language interpreter coordination cue",
            "Federated Dhvani — SSO for other UN agencies (stretch)",
        ],
        metric="3,000 MAU. Under $2/user/month. 99.5% availability. First Member State requests access.",
    )

    build_demo_slide(prs, 11, total)
    build_ui_polish_slide(prs, 12, total)
    build_priority_slide(prs, 13, total)
    build_metrics_slide(prs, 14, total)
    build_risks_slide(prs, 15, total)
    build_ask_slide(prs, 16, total)
    build_next14_slide(prs, 17, total)
    build_closing_slide(prs, 18, total)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "presentation"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "Dhvani_Roadmap.pptx"
    prs.save(str(out_path))
    print(f"wrote {out_path}")


if __name__ == "__main__":
    main()
